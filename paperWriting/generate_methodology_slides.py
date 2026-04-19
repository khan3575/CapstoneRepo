"""
Generate 5 methodology slides (Slides 10-14) for the GNN Brain Tumour Segmentation presentation.
Matches content from PRESENTATION_SPEECH.md — Rifa Sanjida's section.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Color palette ──────────────────────────────────────────────────────────────
BG_DARK       = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
ACCENT_BLUE   = RGBColor(0x1B, 0x6C, 0xA8)   # steel blue
ACCENT_TEAL   = RGBColor(0x00, 0xB4, 0xD8)   # bright teal
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xCC, 0xD6, 0xE0)
YELLOW        = RGBColor(0xFF, 0xD1, 0x66)
GREEN         = RGBColor(0x06, 0xD6, 0xA0)
CARD_BG       = RGBColor(0x1A, 0x2E, 0x44)   # slightly lighter navy for cards
ORANGE        = RGBColor(0xFF, 0x8C, 0x42)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def fill_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill_color, line_color=None, line_width=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        l, t, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             font_size=Pt(18), bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def add_slide_header(slide, slide_num, title, subtitle=None):
    """Dark header bar with slide number + title."""
    fill_bg(slide, BG_DARK)
    # top accent bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.07), ACCENT_TEAL)
    # slide number pill
    add_rect(slide, Inches(0.3), Inches(0.15), Inches(0.45), Inches(0.35), ACCENT_BLUE)
    add_text(slide, str(slide_num),
             Inches(0.3), Inches(0.13), Inches(0.45), Inches(0.4),
             font_size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # title
    add_text(slide, title,
             Inches(0.85), Inches(0.1), Inches(11.5), Inches(0.55),
             font_size=Pt(28), bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.85), Inches(0.55), Inches(11.5), Inches(0.35),
                 font_size=Pt(14), bold=False, color=ACCENT_TEAL)
    # bottom thin line
    add_rect(slide, 0, Inches(7.43), SLIDE_W, Inches(0.07), ACCENT_BLUE)


def add_card(slide, l, t, w, h, title, lines, title_color=ACCENT_TEAL,
             title_size=Pt(15), body_size=Pt(13)):
    """Rounded-look card (rectangle + title + bullet lines)."""
    add_rect(slide, l, t, w, h, CARD_BG, ACCENT_BLUE, line_width=Pt(1.2))
    # title
    add_text(slide, title, l + Inches(0.12), t + Inches(0.08),
             w - Inches(0.2), Inches(0.35),
             font_size=title_size, bold=True, color=title_color)
    # body lines
    body_text = "\n".join(lines)
    add_text(slide, body_text,
             l + Inches(0.12), t + Inches(0.42),
             w - Inches(0.2), h - Inches(0.52),
             font_size=body_size, color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — SLIC SUPERPIXELS
# ══════════════════════════════════════════════════════════════════════════════
def slide_slic(prs):
    slide = blank_slide(prs)
    add_slide_header(slide, 10, "SLIC Superpixels",
                     "Step 1 — Spatial Compression of MRI Slices")

    # ── Left column: algorithm parameters ──────────────────────────────────
    add_card(slide,
             Inches(0.3), Inches(1.0), Inches(4.0), Inches(2.6),
             "Algorithm Parameters",
             ["Algorithm :  SLIC (Simple Linear Iterative Clustering)",
              "Applied on :  T1CE channel only",
              "  (highest tumour contrast)",
              "Target K :  200 superpixels per slice",
              "Compactness :  0.1",
              "Sigma :  0.3",
              "Realised :  ~46 superpixels per slice"])

    # ── Left column: why T1CE ──────────────────────────────────────────────
    add_card(slide,
             Inches(0.3), Inches(3.75), Inches(4.0), Inches(1.7),
             "Why T1CE?",
             ["Contrast agent highlights areas where",
              "blood-brain barrier is broken",
              "→ Active tumour tissue appears bright",
              "→ Clearest tumour boundary signal"])

    # ── Right column: compression visual ──────────────────────────────────
    # Big number highlight
    add_rect(slide, Inches(4.7), Inches(1.0), Inches(8.3), Inches(1.4), CARD_BG,
             ACCENT_TEAL, line_width=Pt(1.5))
    add_text(slide, "2,284×  Spatial Compression",
             Inches(4.9), Inches(1.05), Inches(8.0), Inches(0.5),
             font_size=Pt(22), bold=True, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)
    add_text(slide, "8,928,000 voxels per patient   →   ~3,909 superpixel nodes",
             Inches(4.9), Inches(1.55), Inches(8.0), Inches(0.45),
             font_size=Pt(14), color=WHITE, align=PP_ALIGN.CENTER)

    # Per-slice breakdown
    add_card(slide,
             Inches(4.7), Inches(2.55), Inches(3.9), Inches(2.9),
             "Per Slice",
             ["Raw MRI slice :",
              "  155 × 240 = 37,200 pixels",
              "",
              "After SLIC :",
              "  ~46 superpixels",
              "",
              "Compression per slice :",
              "  ~808×"])

    add_card(slide,
             Inches(8.8), Inches(2.55), Inches(4.2), Inches(2.9),
             "Per Patient (Full Volume)",
             ["MRI volume :",
              "  155 slices × 240 × 240 = 8.9M voxels",
              "",
              "After SLIC :",
              "  ~46 nodes × 155 slices ≈ 7,130 nodes",
              "  (paired-slice → ~3,909 graph nodes)",
              "",
              "This is WHY inference is fast"])

    # bottom note
    add_text(slide,
             "Each superpixel is a compact, semantically meaningful region — far better than raw pixels as a unit of analysis.",
             Inches(0.3), Inches(5.6), Inches(12.7), Inches(0.5),
             font_size=Pt(12), color=LIGHT_GRAY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — GRAPH CONSTRUCTION
# ══════════════════════════════════════════════════════════════════════════════
def slide_graph_construction(prs):
    slide = blank_slide(prs)
    add_slide_header(slide, 11, "Graph Construction",
                     "Step 2 — Building the Paired-Slice Graph")

    # ── Overview card ──────────────────────────────────────────────────────
    add_card(slide,
             Inches(0.3), Inches(1.0), Inches(5.5), Inches(1.4),
             "Graph Unit",
             ["Input :  2 consecutive axial MRI slices",
              "Nodes :  ~92 superpixels   |   Edges :  ~180 connections",
              "One graph unit per paired slice through the full volume"])

    # ── Edge type cards ────────────────────────────────────────────────────
    add_card(slide,
             Inches(0.3), Inches(2.6), Inches(6.1), Inches(2.4),
             "Edge Type 1 — Intra-Slice  (within same slice)",
             ["Method :  Region Adjacency Graph (RAG)",
              "Rule :  4-connectivity boundary adjacency",
              "Connects superpixels that share a physical border",
              "",
              "Captures :  tissue boundaries, local spatial structure",
              "→ If two regions are adjacent, they likely have",
              "   biologically meaningful relationship"],
             title_color=ACCENT_TEAL)

    add_card(slide,
             Inches(6.7), Inches(2.6), Inches(6.3), Inches(2.4),
             "Edge Type 2 — Inter-Slice  (across two slices)",
             ["Method :  k-Nearest Neighbours  (k = 3)",
              "Criteria — edge only if BOTH hold:",
              "  • IoU overlap  >  0.1  (10%)",
              "  • Centroid distance  <  10 mm",
              "",
              "Captures :  3D spatial continuity of tumour",
              "→ Links the same anatomical region across",
              "   consecutive slices"],
             title_color=YELLOW)

    # ── Stats bar ──────────────────────────────────────────────────────────
    add_rect(slide, Inches(0.3), Inches(5.2), Inches(12.7), Inches(0.95), CARD_BG,
             ACCENT_BLUE, line_width=Pt(1))

    stats = [
        ("~92", "nodes per graph"),
        ("~180", "edges per graph"),
        ("2", "edge types"),
        ("2,284×", "compression vs raw voxels"),
    ]
    col_w = Inches(3.1)
    for i, (num, label) in enumerate(stats):
        x = Inches(0.5) + i * col_w
        add_text(slide, num,  x, Inches(5.22), col_w, Inches(0.42),
                 font_size=Pt(22), bold=True, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)
        add_text(slide, label, x, Inches(5.63), col_w, Inches(0.3),
                 font_size=Pt(11), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    add_text(slide,
             "Two-slice structure lets the model learn both spatial adjacency (within slice) and volumetric continuity (across slices).",
             Inches(0.3), Inches(6.3), Inches(12.7), Inches(0.45),
             font_size=Pt(12), color=LIGHT_GRAY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 12 — NODE FEATURES
# ══════════════════════════════════════════════════════════════════════════════
def slide_node_features(prs):
    slide = blank_slide(prs)
    add_slide_header(slide, 12, "Node Features",
                     "Step 3 — 15-Dimensional Feature Vector per Superpixel")

    # ── Header label ───────────────────────────────────────────────────────
    add_rect(slide, Inches(0.3), Inches(0.95), Inches(12.7), Inches(0.45),
             ACCENT_BLUE)
    add_text(slide, "Each node (superpixel) is described by 15 features across 3 groups",
             Inches(0.4), Inches(0.97), Inches(12.5), Inches(0.4),
             font_size=Pt(14), bold=False, color=WHITE, align=PP_ALIGN.CENTER)

    # ── Group 1: Intensity ─────────────────────────────────────────────────
    add_card(slide,
             Inches(0.3), Inches(1.55), Inches(4.1), Inches(4.3),
             "Group 1 — Intensity Features  (8)",
             ["For each of the 4 MRI modalities:",
              "  T1 · T1CE · T2 · FLAIR",
              "",
              "Extract 2 statistics:",
              "  • Mean intensity of pixels in superpixel",
              "  • Standard deviation of intensity",
              "",
              "4 modalities × 2 stats = 8 features",
              "",
              "Tells the model WHAT the tissue looks like",
              "in each imaging contrast"],
             title_color=ACCENT_TEAL,
             title_size=Pt(13), body_size=Pt(12))

    # ── Group 2: Spatial ───────────────────────────────────────────────────
    add_card(slide,
             Inches(4.6), Inches(1.55), Inches(3.9), Inches(4.3),
             "Group 2 — Spatial Features  (4)",
             ["1.  Area",
              "    Pixel count inside superpixel",
              "",
              "2.  Normalised Area",
              "    Area ÷ total slice area",
              "",
              "3.  Y-Centroid",
              "    Vertical position in slice",
              "",
              "4.  X-Centroid",
              "    Horizontal position in slice",
              "",
              "Tells the model WHERE the region is"],
             title_color=YELLOW,
             title_size=Pt(13), body_size=Pt(12))

    # ── Group 3: Morphological ─────────────────────────────────────────────
    add_card(slide,
             Inches(8.7), Inches(1.55), Inches(4.3), Inches(4.3),
             "Group 3 — Morphological Features  (3)",
             ["1.  Perimeter",
              "    Boundary length of superpixel",
              "",
              "2.  Compactness",
              "    Perimeter² ÷ (4π × Area)",
              "    Circle = 1.0, irregular shapes < 1",
              "",
              "3.  Intensity Range",
              "    Max − Min pixel intensity",
              "",
              "Tells the model the SHAPE and",
              "texture complexity of the region"],
             title_color=ORANGE,
             title_size=Pt(13), body_size=Pt(12))

    # ── Summary bar ────────────────────────────────────────────────────────
    add_rect(slide, Inches(0.3), Inches(6.0), Inches(12.7), Inches(0.55),
             CARD_BG, ACCENT_BLUE, line_width=Pt(1))
    add_text(slide,
             "Total:  8 (intensity)  +  4 (spatial)  +  3 (morphological)  =  15 features per node",
             Inches(0.5), Inches(6.07), Inches(12.3), Inches(0.42),
             font_size=Pt(15), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide,
             "Combined, these features capture what the tissue is, where it is, and what shape it has — sufficient to classify tumour vs. healthy.",
             Inches(0.3), Inches(6.68), Inches(12.7), Inches(0.45),
             font_size=Pt(12), color=LIGHT_GRAY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 13 — GRAPHSAGE ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
def slide_graphsage(prs):
    slide = blank_slide(prs)
    add_slide_header(slide, 13, "GraphSAGE Architecture",
                     "Graph Sample and Aggregate — Inductive GNN")

    # ── Why GraphSAGE card ─────────────────────────────────────────────────
    add_card(slide,
             Inches(0.3), Inches(1.0), Inches(6.2), Inches(2.5),
             "Why GraphSAGE? (Not GCN or GAT)",
             ["Key property:  INDUCTIVE learning",
              "",
              "GCN / GAT are transductive:",
              "  → require the full graph at training time",
              "  → cannot generalise to unseen graphs",
              "",
              "GraphSAGE learns a neighbourhood SAMPLING",
              "FUNCTION, not a fixed adjacency matrix",
              "  → works on any graph at inference time",
              "  → enables zero-shot transfer to BraTS 2023"],
             title_color=ACCENT_TEAL)

    # ── Architecture specs ─────────────────────────────────────────────────
    add_card(slide,
             Inches(6.8), Inches(1.0), Inches(6.2), Inches(2.5),
             "Architecture Specifications",
             ["Layers :            5  GraphSAGE layers",
              "Hidden channels :   256  per layer",
              "Output dimension :  64",
              "Aggregation :       Mean pooling",
              "Total parameters :  439,041",
              "                    (< 0.5 million)",
              "Task :              Binary node classification",
              "                    (tumour vs. healthy per superpixel)"],
             title_color=YELLOW)

    # ── Layer flow diagram (text-based) ────────────────────────────────────
    add_rect(slide, Inches(0.3), Inches(3.7), Inches(12.7), Inches(1.55), CARD_BG,
             ACCENT_BLUE, line_width=Pt(1))
    add_text(slide, "Message Passing Flow  (per layer):",
             Inches(0.5), Inches(3.75), Inches(4.0), Inches(0.35),
             font_size=Pt(13), bold=True, color=ACCENT_TEAL)

    flow_steps = [
        ("Input\n15-dim", ACCENT_BLUE),
        ("Layer 1\n256-dim", CARD_BG),
        ("Layer 2\n256-dim", CARD_BG),
        ("Layer 3\n256-dim", CARD_BG),
        ("Layer 4\n256-dim", CARD_BG),
        ("Layer 5\n64-dim", ACCENT_BLUE),
        ("Sigmoid\n0 or 1", GREEN),
    ]
    box_w = Inches(1.6)
    start_x = Inches(0.45)
    box_y = Inches(4.15)
    for i, (label, color) in enumerate(flow_steps):
        x = start_x + i * Inches(1.82)
        add_rect(slide, x, box_y, box_w, Inches(0.8), color,
                 ACCENT_TEAL, line_width=Pt(1))
        add_text(slide, label, x, box_y + Inches(0.05), box_w, Inches(0.75),
                 font_size=Pt(11), bold=(i in [0, 5, 6]), color=WHITE,
                 align=PP_ALIGN.CENTER)
        if i < len(flow_steps) - 1:
            arrow_x = x + box_w + Inches(0.02)
            add_text(slide, "→", arrow_x, box_y + Inches(0.25), Inches(0.18), Inches(0.35),
                     font_size=Pt(16), bold=True, color=ACCENT_TEAL,
                     align=PP_ALIGN.CENTER)

    # ── Aggregation explanation ────────────────────────────────────────────
    add_card(slide,
             Inches(0.3), Inches(5.4), Inches(6.1), Inches(1.6),
             "Mean Aggregation (per layer)",
             ["1.  For each node, collect features from neighbours",
              "2.  Average all neighbour feature vectors",
              "3.  Concatenate with own features",
              "4.  Apply linear transform + activation (ReLU)"],
             body_size=Pt(12))

    add_card(slide,
             Inches(6.65), Inches(5.4), Inches(6.35), Inches(1.6),
             "Why 5 Layers? (Ablation-Justified)",
             ["6 layers → same Dice (84.00%) but 130K more params",
              "512-dim  → +4.75pp Dice but 4× more parameters",
              "GAT      → worse Dice, 2.7× more parameters",
              "→ 5L, 256-dim is the optimal efficiency point"],
             body_size=Pt(12))


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 14 — TRAINING PROTOCOL
# ══════════════════════════════════════════════════════════════════════════════
def slide_training(prs):
    slide = blank_slide(prs)
    add_slide_header(slide, 14, "Training Protocol",
                     "Optimisation, Loss, and Cross-Validation Setup")

    # ── Top row: 3 cards ───────────────────────────────────────────────────
    add_card(slide,
             Inches(0.3), Inches(1.0), Inches(4.1), Inches(2.6),
             "Optimiser & Scheduler",
             ["Optimiser :  AdamW",
              "  Learning rate :  1e-3",
              "  Weight decay :  0.01",
              "",
              "Scheduler :  OneCycleLR",
              "  Warmup :  30% of total steps",
              "  (gradual LR ramp-up to prevent",
              "   early training instability)"],
             title_color=ACCENT_TEAL)

    add_card(slide,
             Inches(4.6), Inches(1.0), Inches(4.1), Inches(2.6),
             "Batch & Precision",
             ["Batch size :  24 graphs",
              "Grad. accumulation :  ×2",
              "Effective batch :  48 graphs",
              "",
              "Precision :  AMP FP16",
              "  (mixed precision — halves",
              "   GPU memory, speeds training)",
              "",
              "Max epochs :  50",
              "Early stop patience :  10  (never triggered)"],
             title_color=YELLOW)

    add_card(slide,
             Inches(8.9), Inches(1.0), Inches(4.1), Inches(2.6),
             "Loss Function",
             ["Loss :  BCEWithLogitsLoss",
              "",
              "Class imbalance problem:",
              "  ~19 healthy : 1 tumour superpixel",
              "",
              "Solution: positive class weight",
              "  w+ = 9.0",
              "  Forces model to penalise missed",
              "  tumour detections more heavily",
              "",
              "Dice smoothing ε = 1e-7"],
             title_color=ORANGE)

    # ── Cross-validation card ──────────────────────────────────────────────
    add_card(slide,
             Inches(0.3), Inches(3.75), Inches(6.3), Inches(2.35),
             "5-Fold Cross-Validation",
             ["1,000 BraTS 2021 patients split at PATIENT level:",
              "",
              "  Each fold:   720 train  /  80 val  /  200 test",
              "",
              "5 models trained independently on 5 splits",
              "",
              "Ensemble:  average sigmoid probabilities",
              "           from all 5 folds (soft voting)",
              "Threshold: τ = 0.5"],
             title_color=GREEN)

    # ── Hardware card ──────────────────────────────────────────────────────
    add_card(slide,
             Inches(6.85), Inches(3.75), Inches(6.15), Inches(2.35),
             "Hardware (Consumer-Grade Only)",
             ["CPU :   Intel Core i7-10700",
              "RAM :   32 GB",
              "GPU :   NVIDIA RTX 2060  —  6 GB VRAM",
              "",
              "Python 3.12  ·  PyTorch 2.8.0  (CUDA 12.8)",
              "PyTorch Geometric 2.6.1",
              "scikit-image 0.25.2  ·  NumPy 2.3.3",
              "",
              "Everything trained and tested on 6 GB VRAM"],
             title_color=LIGHT_GRAY)

    add_text(slide,
             "Entire training pipeline — from superpixel construction to ensemble inference — runs on hardware available in most university labs.",
             Inches(0.3), Inches(6.25), Inches(12.7), Inches(0.45),
             font_size=Pt(12), color=LIGHT_GRAY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════════════════════
if __name__ == "__main__":
    prs = new_prs()

    slide_slic(prs)
    slide_graph_construction(prs)
    slide_node_features(prs)
    slide_graphsage(prs)
    slide_training(prs)

    out = "/mnt/bigdata/capstone/brats_gnn_segmentation/paperWriting/methodology_slides_10_14.pptx"
    prs.save(out)
    print(f"Saved → {out}")
