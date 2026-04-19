"""
Generate compressed 15-slide presentation for 10-minute defence.
GNN Brain Tumour Segmentation — BUBT Capstone 2026
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ────────────────────────────────────────────────────────────────────
BG_DARK     = RGBColor(0x0D, 0x1B, 0x2A)
ACCENT_BLUE = RGBColor(0x1B, 0x6C, 0xA8)
ACCENT_TEAL = RGBColor(0x00, 0xB4, 0xD8)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xCC, 0xD6, 0xE0)
YELLOW      = RGBColor(0xFF, 0xD1, 0x66)
GREEN       = RGBColor(0x06, 0xD6, 0xA0)
CARD_BG     = RGBColor(0x1A, 0x2E, 0x44)
ORANGE      = RGBColor(0xFF, 0x8C, 0x42)
RED         = RGBColor(0xEF, 0x47, 0x6F)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill_color, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             font_size=Pt(18), bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def add_card(slide, l, t, w, h, title, lines,
             title_color=ACCENT_TEAL, title_size=Pt(15), body_size=Pt(13)):
    add_rect(slide, l, t, w, h, CARD_BG, ACCENT_BLUE, line_width=Pt(1.2))
    add_text(slide, title,
             l + Inches(0.12), t + Inches(0.08),
             w - Inches(0.2), Inches(0.35),
             font_size=title_size, bold=True, color=title_color)
    add_text(slide, "\n".join(lines),
             l + Inches(0.12), t + Inches(0.42),
             w - Inches(0.2), h - Inches(0.52),
             font_size=body_size, color=LIGHT_GRAY)


def add_header(slide, num, title, subtitle=None):
    fill_bg(slide, BG_DARK)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.07), ACCENT_TEAL)
    add_rect(slide, Inches(0.3), Inches(0.15), Inches(0.45), Inches(0.35), ACCENT_BLUE)
    add_text(slide, str(num),
             Inches(0.3), Inches(0.13), Inches(0.45), Inches(0.4),
             font_size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title,
             Inches(0.85), Inches(0.1), Inches(11.5), Inches(0.55),
             font_size=Pt(28), bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.85), Inches(0.55), Inches(11.5), Inches(0.35),
                 font_size=Pt(14), color=ACCENT_TEAL)
    add_rect(slide, 0, Inches(7.43), SLIDE_W, Inches(0.07), ACCENT_BLUE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE + TEAM
# ══════════════════════════════════════════════════════════════════════════════
def slide_title(prs):
    slide = blank_slide(prs)
    fill_bg(slide, BG_DARK)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), ACCENT_TEAL)
    add_rect(slide, 0, Inches(7.42), SLIDE_W, Inches(0.08), ACCENT_TEAL)

    # Centre title block
    add_rect(slide, Inches(0.8), Inches(0.9), Inches(11.73), Inches(2.2), CARD_BG,
             ACCENT_TEAL, line_width=Pt(1.5))
    add_text(slide,
             "Efficient Brain Tumour Segmentation",
             Inches(0.9), Inches(1.0), Inches(11.5), Inches(0.7),
             font_size=Pt(30), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide,
             "using Graph Neural Networks on BraTS Datasets",
             Inches(0.9), Inches(1.68), Inches(11.5), Inches(0.55),
             font_size=Pt(22), bold=False, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)

    # Team members
    add_rect(slide, Inches(0.8), Inches(3.3), Inches(5.5), Inches(3.3), CARD_BG,
             ACCENT_BLUE, line_width=Pt(1))
    add_text(slide, "Team Members",
             Inches(0.95), Inches(3.38), Inches(5.2), Inches(0.38),
             font_size=Pt(14), bold=True, color=ACCENT_TEAL)
    members = [
        "Sakib Khan",
        "Rifa Sanjida",
        "Kishor Kumar Das",
        "Md. Mahamudul Hasan",
        "Md. Minhajur Rahman",
    ]
    for i, name in enumerate(members):
        add_text(slide, f"  {i+1}.  {name}",
                 Inches(0.95), Inches(3.82) + i * Inches(0.48),
                 Inches(5.2), Inches(0.45),
                 font_size=Pt(14), color=LIGHT_GRAY)

    # Supervisor + institution
    add_rect(slide, Inches(6.8), Inches(3.3), Inches(6.1), Inches(3.3), CARD_BG,
             ACCENT_BLUE, line_width=Pt(1))
    add_text(slide, "Supervised by",
             Inches(6.95), Inches(3.38), Inches(5.8), Inches(0.38),
             font_size=Pt(14), bold=True, color=ACCENT_TEAL)
    add_text(slide, "Mr. Shamim Ahmed",
             Inches(6.95), Inches(3.82), Inches(5.8), Inches(0.45),
             font_size=Pt(16), bold=True, color=WHITE)
    add_text(slide, "Department of Computer Science & Engineering\nBangladesh University of Business & Technology  (BUBT)",
             Inches(6.95), Inches(4.32), Inches(5.8), Inches(0.9),
             font_size=Pt(13), color=LIGHT_GRAY)
    add_text(slide, "Capstone Project  ·  2026",
             Inches(6.95), Inches(5.35), Inches(5.8), Inches(0.38),
             font_size=Pt(13), color=YELLOW)
    add_text(slide, "BraTS 2021  ·  BraTS 2023  ·  GraphSAGE  ·  Python / PyTorch",
             Inches(6.95), Inches(5.82), Inches(5.8), Inches(0.38),
             font_size=Pt(12), color=LIGHT_GRAY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM STATEMENT
# ══════════════════════════════════════════════════════════════════════════════
def slide_problem(prs):
    slide = blank_slide(prs)
    add_header(slide, 2, "Problem Statement",
               "Why existing methods fall short in clinical settings")

    add_card(slide, Inches(0.3), Inches(1.0), Inches(6.1), Inches(2.8),
             "The Clinical Need",
             ["Brain tumour segmentation from MRI is essential for:",
              "  • Surgical planning and resection guidance",
              "  • Radiotherapy targeting",
              "  • Treatment response monitoring",
              "",
              "Multi-modal MRI (T1, T1CE, T2, FLAIR) provides",
              "complementary tissue contrast — manual annotation",
              "is slow, expert-dependent, and not reproducible."],
             title_color=ACCENT_TEAL)

    add_card(slide, Inches(6.7), Inches(1.0), Inches(6.3), Inches(2.8),
             "Problem with Existing Deep Learning Methods",
             ["CNNs / U-Net variants achieve high Dice scores BUT:",
              "",
              "  • Require 16–32 GB VRAM",
              "  • 69M–100M+ parameters",
              "  • Long inference times (10+ seconds)",
              "  • Not deployable in resource-constrained",
              "    clinical settings or low-income countries",
              "",
              "  GPU-heavy pipelines ≠ accessible medicine"],
             title_color=RED)

    add_rect(slide, Inches(0.3), Inches(4.0), Inches(12.7), Inches(1.1),
             CARD_BG, ACCENT_TEAL, line_width=Pt(1.5))
    add_text(slide, "Our Proposal",
             Inches(0.5), Inches(4.06), Inches(3.0), Inches(0.38),
             font_size=Pt(14), bold=True, color=ACCENT_TEAL)
    add_text(slide,
             "Replace pixel-level CNN with a lightweight Graph Neural Network on SLIC superpixels.\n"
             "Same accuracy — fraction of the compute.",
             Inches(0.5), Inches(4.46), Inches(12.3), Inches(0.6),
             font_size=Pt(15), color=WHITE)

    add_card(slide, Inches(0.3), Inches(5.25), Inches(4.0), Inches(1.75),
             "Class Imbalance",
             ["~19 : 1  (healthy vs. tumour superpixels)",
              "Addressed via positive class weight  w+ = 9.0"],
             title_color=YELLOW, title_size=Pt(13), body_size=Pt(12))

    add_card(slide, Inches(4.5), Inches(5.25), Inches(4.0), Inches(1.75),
             "Hardware Constraint",
             ["All experiments on  NVIDIA RTX 2060  (6 GB VRAM)",
              "Consumer-grade GPU — no HPC cluster needed"],
             title_color=YELLOW, title_size=Pt(13), body_size=Pt(12))

    add_card(slide, Inches(8.7), Inches(5.25), Inches(4.3), Inches(1.75),
             "Generalisation Goal",
             ["Train on BraTS 2021 → test zero-shot on BraTS 2023",
              "Different scanner protocols, no retraining"],
             title_color=YELLOW, title_size=Pt(13), body_size=Pt(12))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — RESEARCH OBJECTIVES
# ══════════════════════════════════════════════════════════════════════════════
def slide_objectives(prs):
    slide = blank_slide(prs)
    add_header(slide, 3, "Research Objectives",
               "Four measurable goals driving this work")

    objectives = [
        ("01", "GNN Segmentation Pipeline",
         "Develop a complete GNN-based segmentation pipeline using SLIC superpixel graphs "
         "on multi-modal BraTS MRI (T1, T1CE, T2, FLAIR).",
         ACCENT_TEAL),
        ("02", "Competitive Dice Score",
         "Achieve Dice score comparable to or better than CNN / U-Net baselines "
         "on the BraTS 2021 benchmark.",
         YELLOW),
        ("03", "Dramatically Reduce Compute",
         "Reduce GPU memory, inference time, and parameter count by orders of magnitude "
         "vs. standard deep learning methods — targeting 6 GB VRAM deployment.",
         ORANGE),
        ("04", "Cross-Dataset Generalisation",
         "Demonstrate zero-shot transfer from BraTS 2021 to BraTS 2023 without any retraining, "
         "validating the inductive property of GraphSAGE.",
         GREEN),
    ]

    for i, (num, title, body, color) in enumerate(objectives):
        col = i % 2
        row = i // 2
        x = Inches(0.3) + col * Inches(6.55)
        y = Inches(1.05) + row * Inches(2.85)
        w = Inches(6.3)
        h = Inches(2.6)
        add_rect(slide, x, y, w, h, CARD_BG, color, line_width=Pt(1.5))
        # Big number
        add_text(slide, num,
                 x + Inches(0.12), y + Inches(0.08),
                 Inches(0.65), Inches(0.65),
                 font_size=Pt(28), bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, title,
                 x + Inches(0.85), y + Inches(0.1),
                 w - Inches(1.0), Inches(0.45),
                 font_size=Pt(16), bold=True, color=WHITE)
        add_text(slide, body,
                 x + Inches(0.12), y + Inches(0.65),
                 w - Inches(0.25), h - Inches(0.75),
                 font_size=Pt(13), color=LIGHT_GRAY)

    add_text(slide,
             "Each objective maps to a measurable result — all four are met by the end of this presentation.",
             Inches(0.3), Inches(6.82), Inches(12.7), Inches(0.4),
             font_size=Pt(12), color=LIGHT_GRAY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — DATASET
# ══════════════════════════════════════════════════════════════════════════════
def slide_dataset(prs):
    slide = blank_slide(prs)
    add_header(slide, 4, "Dataset",
               "BraTS 2021 (primary) + BraTS 2023 (zero-shot generalisation test)")

    add_card(slide, Inches(0.3), Inches(1.0), Inches(4.0), Inches(5.55),
             "BraTS 2021 — Primary",
             ["1,251 glioma patients total",
              "",
              "Training / Evaluation pool:",
              "  1,000 patients",
              "  → 5-fold cross-validation",
              "  → 720 train / 80 val / 200 test per fold",
              "",
              "Sealed held-out set:",
              "  251 patients",
              "  → never touched during training",
              "  → final ensemble evaluated here",
              "",
              "4 MRI modalities per patient:",
              "  T1   T1CE   T2   FLAIR",
              "",
              "Label:  Whole Tumour binary mask",
              "  (tumour  vs.  healthy)"],
             title_color=ACCENT_TEAL, body_size=Pt(12))

    add_card(slide, Inches(4.5), Inches(1.0), Inches(4.0), Inches(5.55),
             "BraTS 2023 — Zero-Shot Test",
             ["Used ONLY for generalisation test",
              "No retraining, no fine-tuning",
              "",
              "Different acquisition protocols",
              "vs. BraTS 2021:",
              "  • Different scanner manufacturers",
              "  • Different site calibrations",
              "  • Harder distribution shift",
              "",
              "Same label schema as BraTS 2021",
              "→ direct Dice comparison possible",
              "",
              "Tests whether GraphSAGE's",
              "inductive property holds in the",
              "real world"],
             title_color=GREEN, body_size=Pt(12))

    add_card(slide, Inches(8.7), Inches(1.0), Inches(4.3), Inches(5.55),
             "Data Characteristics",
             ["Modality contrast roles:",
              "  T1CE  →  active tumour (bright ring)",
              "  FLAIR →  oedema / infiltration",
              "  T2    →  water content",
              "  T1    →  anatomy baseline",
              "",
              "SLIC applied on T1CE only",
              "  (highest tumour contrast)",
              "",
              "Superpixel-level class balance:",
              "  ~19 healthy : 1 tumour",
              "  (19 : 1 imbalance)",
              "",
              "Positive class weight  w+ = 9.0",
              "  to compensate imbalance",
              "",
              "Spatial compression:",
              "  8.9M voxels → ~3,909 nodes"],
             title_color=YELLOW, body_size=Pt(12))

    add_text(slide,
             "All splits are patient-level — no data leakage between train and test sets.",
             Inches(0.3), Inches(6.7), Inches(12.7), Inches(0.45),
             font_size=Pt(12), color=LIGHT_GRAY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — METHODOLOGY OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
def slide_methodology_overview(prs):
    slide = blank_slide(prs)
    add_header(slide, 5, "Proposed Methodology",
               "Five-stage pipeline: MRI → Superpixels → Graph → GNN → Segmentation")

    # Pipeline boxes
    pipeline = [
        ("MRI Input\n4 Modalities", ACCENT_BLUE),
        ("SLIC\nSuperpixels", ACCENT_TEAL),
        ("Graph\nConstruction", YELLOW),
        ("GraphSAGE\n5 Layers", ORANGE),
        ("Binary\nSegmentation", GREEN),
    ]
    bw = Inches(2.2)
    bh = Inches(1.4)
    gap = Inches(0.22)
    start_x = Inches(0.35)
    by = Inches(1.15)
    for i, (label, color) in enumerate(pipeline):
        x = start_x + i * (bw + gap)
        add_rect(slide, x, by, bw, bh, color, WHITE, line_width=Pt(1))
        add_text(slide, label, x, by + Inches(0.3), bw, Inches(0.85),
                 font_size=Pt(15), bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
        if i < len(pipeline) - 1:
            ax = x + bw + Inches(0.03)
            add_text(slide, "→", ax, by + Inches(0.45), gap, Inches(0.5),
                     font_size=Pt(20), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Detail cards for each stage
    details = [
        ("MRI Input",
         ["T1, T1CE, T2, FLAIR", "155 axial slices", "240×240 per slice", "4-channel input"]),
        ("SLIC Superpixels",
         ["Applied on T1CE", "K=200 target", "~46 per slice", "2,284× compression"]),
        ("Graph Construction",
         ["Paired-slice (2 slices)", "~92 nodes, ~180 edges", "Intra: RAG adjacency", "Inter: kNN  k=3"]),
        ("GraphSAGE GNN",
         ["5 layers, 256 hidden", "64-dim output", "439,041 params", "Mean aggregation"]),
        ("Segmentation",
         ["Node classification", "Tumour / Healthy", "5-fold ensemble", "τ = 0.5"]),
    ]
    bh2 = Inches(2.5)
    by2 = Inches(2.75)
    for i, (title, lines) in enumerate(details):
        x = start_x + i * (bw + gap)
        add_card(slide, x, by2, bw, bh2, title, lines,
                 title_color=ACCENT_TEAL, title_size=Pt(12), body_size=Pt(11))

    # Key insight bar
    add_rect(slide, Inches(0.3), Inches(5.45), Inches(12.7), Inches(0.9),
             CARD_BG, ACCENT_TEAL, line_width=Pt(1.2))
    add_text(slide, "Key Insight:",
             Inches(0.5), Inches(5.52), Inches(1.5), Inches(0.35),
             font_size=Pt(13), bold=True, color=ACCENT_TEAL)
    add_text(slide,
             "By operating on ~92 superpixel nodes rather than 37,200 raw pixels per slice, the GNN "
             "achieves 2,284× spatial compression — enabling inference in 75 ms on a 6 GB consumer GPU.",
             Inches(2.1), Inches(5.52), Inches(10.7), Inches(0.75),
             font_size=Pt(13), color=WHITE)

    add_text(slide,
             "The following slides detail each stage.",
             Inches(0.3), Inches(6.5), Inches(12.7), Inches(0.4),
             font_size=Pt(12), color=LIGHT_GRAY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — SLIC SUPERPIXELS
# ══════════════════════════════════════════════════════════════════════════════
def slide_slic(prs):
    slide = blank_slide(prs)
    add_header(slide, 6, "SLIC Superpixels",
               "Step 1 — Spatial Compression of MRI Slices")

    add_card(slide, Inches(0.3), Inches(1.0), Inches(4.0), Inches(2.6),
             "Algorithm Parameters",
             ["Algorithm :  SLIC (Simple Linear Iterative Clustering)",
              "Applied on :  T1CE channel only",
              "  (highest tumour contrast)",
              "Target K :  200 superpixels per slice",
              "Compactness :  0.1",
              "Sigma :  0.3",
              "Realised :  ~46 superpixels per slice"])

    add_card(slide, Inches(0.3), Inches(3.75), Inches(4.0), Inches(1.7),
             "Why T1CE?",
             ["Contrast agent highlights areas where blood-brain barrier is broken",
              "→ Active tumour tissue appears bright",
              "→ Clearest tumour boundary signal"])

    add_rect(slide, Inches(4.7), Inches(1.0), Inches(8.3), Inches(1.4), CARD_BG,
             ACCENT_TEAL, line_width=Pt(1.5))
    add_text(slide, "2,284×  Spatial Compression",
             Inches(4.9), Inches(1.05), Inches(8.0), Inches(0.5),
             font_size=Pt(22), bold=True, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)
    add_text(slide, "8,928,000 voxels per patient   →   ~3,909 superpixel nodes",
             Inches(4.9), Inches(1.55), Inches(8.0), Inches(0.45),
             font_size=Pt(14), color=WHITE, align=PP_ALIGN.CENTER)

    add_card(slide, Inches(4.7), Inches(2.55), Inches(3.9), Inches(2.9),
             "Per Slice",
             ["Raw MRI slice :",
              "  155 × 240 = 37,200 pixels",
              "",
              "After SLIC :",
              "  ~46 superpixels",
              "",
              "Compression per slice :  ~808×"])

    add_card(slide, Inches(8.8), Inches(2.55), Inches(4.2), Inches(2.9),
             "Per Patient (Full Volume)",
             ["MRI volume :",
              "  155 slices × 240 × 240 = 8.9M voxels",
              "",
              "After SLIC :",
              "  ~46 nodes × 155 slices ≈ 7,130 nodes",
              "  (paired-slice → ~3,909 graph nodes)",
              "",
              "This is WHY inference is fast"])

    add_text(slide,
             "Each superpixel is a compact, semantically meaningful region — far better than raw pixels as a unit of analysis.",
             Inches(0.3), Inches(5.6), Inches(12.7), Inches(0.5),
             font_size=Pt(12), color=LIGHT_GRAY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — GRAPH CONSTRUCTION
# ══════════════════════════════════════════════════════════════════════════════
def slide_graph_construction(prs):
    slide = blank_slide(prs)
    add_header(slide, 7, "Graph Construction",
               "Step 2 — Building the Paired-Slice Graph")

    add_card(slide, Inches(0.3), Inches(1.0), Inches(5.5), Inches(1.4),
             "Graph Unit",
             ["Input :  2 consecutive axial MRI slices",
              "Nodes :  ~92 superpixels   |   Edges :  ~180 connections"])

    add_card(slide, Inches(0.3), Inches(2.6), Inches(6.1), Inches(2.4),
             "Edge Type 1 — Intra-Slice  (within same slice)",
             ["Method :  Region Adjacency Graph (RAG)",
              "Rule :  4-connectivity boundary adjacency",
              "Connects superpixels that share a physical border",
              "",
              "Captures :  tissue boundaries, local spatial structure",
              "→ If two regions are adjacent, they likely share a",
              "   biologically meaningful boundary"],
             title_color=ACCENT_TEAL)

    add_card(slide, Inches(6.7), Inches(2.6), Inches(6.3), Inches(2.4),
             "Edge Type 2 — Inter-Slice  (across two slices)",
             ["Method :  k-Nearest Neighbours  (k = 3)",
              "Criteria — edge only if BOTH hold:",
              "  • IoU overlap  >  0.1  (10%)",
              "  • Centroid distance  <  10 mm",
              "",
              "Captures :  3D spatial continuity of tumour",
              "→ Links same anatomical region across slices"],
             title_color=YELLOW)

    add_rect(slide, Inches(0.3), Inches(5.2), Inches(12.7), Inches(0.95), CARD_BG,
             ACCENT_BLUE, line_width=Pt(1))
    stats = [("~92", "nodes per graph"), ("~180", "edges per graph"),
             ("2", "edge types"), ("2,284×", "compression vs raw")]
    col_w = Inches(3.1)
    for i, (num, label) in enumerate(stats):
        x = Inches(0.5) + i * col_w
        add_text(slide, num, x, Inches(5.22), col_w, Inches(0.42),
                 font_size=Pt(22), bold=True, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)
        add_text(slide, label, x, Inches(5.63), col_w, Inches(0.3),
                 font_size=Pt(11), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    add_text(slide,
             "Two-slice structure lets the model learn both spatial adjacency (within slice) and volumetric continuity (across slices).",
             Inches(0.3), Inches(6.3), Inches(12.7), Inches(0.45),
             font_size=Pt(12), color=LIGHT_GRAY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — NODE FEATURES + GRAPHSAGE (merged)
# ══════════════════════════════════════════════════════════════════════════════
def slide_features_graphsage(prs):
    slide = blank_slide(prs)
    add_header(slide, 8, "Node Features  ·  GraphSAGE Architecture",
               "15-dim feature vector per node  →  5-layer inductive GNN")

    # ── Node features (left half) ─────────────────────────────────────────
    add_card(slide, Inches(0.3), Inches(1.0), Inches(3.9), Inches(5.55),
             "Group 1 — Intensity  (8)",
             ["For each of 4 MRI modalities:",
              "  T1  ·  T1CE  ·  T2  ·  FLAIR",
              "",
              "Extract:  mean + std intensity",
              "4 modalities × 2 = 8 features",
              "",
              "Tells the model WHAT the",
              "tissue looks like"],
             title_color=ACCENT_TEAL, title_size=Pt(13), body_size=Pt(12))

    add_card(slide, Inches(4.4), Inches(1.0), Inches(3.0), Inches(2.65),
             "Group 2 — Spatial  (4)",
             ["1. Area (pixel count)",
              "2. Normalised area",
              "3. Y-centroid",
              "4. X-centroid",
              "",
              "Tells the model WHERE"],
             title_color=YELLOW, title_size=Pt(13), body_size=Pt(12))

    add_card(slide, Inches(4.4), Inches(3.85), Inches(3.0), Inches(2.7),
             "Group 3 — Morphological  (3)",
             ["1. Perimeter",
              "2. Compactness",
              "   (Perimeter² / 4π·Area)",
              "3. Intensity Range",
              "",
              "Tells the model SHAPE"],
             title_color=ORANGE, title_size=Pt(13), body_size=Pt(12))

    # ── GraphSAGE (right half) ─────────────────────────────────────────────
    add_card(slide, Inches(7.65), Inches(1.0), Inches(5.35), Inches(2.65),
             "GraphSAGE — Why Inductive?",
             ["GCN/GAT = transductive (fixed graph at train)",
              "GraphSAGE = inductive (learns aggregation fn)",
              "",
              "→ Works on any unseen graph at inference",
              "→ Enables zero-shot BraTS 2023 transfer",
              "→ No retraining needed for new patients"],
             title_color=ACCENT_TEAL, title_size=Pt(13), body_size=Pt(12))

    add_card(slide, Inches(7.65), Inches(3.85), Inches(5.35), Inches(2.7),
             "Architecture Specs",
             ["Layers :          5  GraphSAGE layers",
              "Hidden channels : 256  per layer",
              "Output dim :      64",
              "Aggregation :     Mean pooling",
              "Total params :    439,041  (< 0.5M)",
              "",
              "Task : Binary node classification",
              "  15-dim → 256 → 256 → 256 → 256 → 64 → σ"],
             title_color=YELLOW, title_size=Pt(13), body_size=Pt(12))

    # Summary bar
    add_rect(slide, Inches(0.3), Inches(6.68), Inches(12.7), Inches(0.45),
             CARD_BG, ACCENT_BLUE, line_width=Pt(1))
    add_text(slide,
             "8 + 4 + 3 = 15 features per node   →   5-layer GraphSAGE   →   439K params  (157× fewer than U-Net)",
             Inches(0.5), Inches(6.73), Inches(12.3), Inches(0.38),
             font_size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — TRAINING PROTOCOL
# ══════════════════════════════════════════════════════════════════════════════
def slide_training(prs):
    slide = blank_slide(prs)
    add_header(slide, 9, "Training Protocol",
               "Optimisation, Loss, and Cross-Validation Setup")

    add_card(slide, Inches(0.3), Inches(1.0), Inches(4.1), Inches(2.6),
             "Optimiser & Scheduler",
             ["Optimiser :  AdamW",
              "  Learning rate :  1e-3",
              "  Weight decay :  0.01",
              "",
              "Scheduler :  OneCycleLR",
              "  Warmup :  30% of total steps",
              "  Gradual LR ramp prevents early",
              "  training instability"])

    add_card(slide, Inches(4.6), Inches(1.0), Inches(4.1), Inches(2.6),
             "Batch & Precision",
             ["Batch size :  24 graphs",
              "Grad. accumulation :  ×2",
              "Effective batch :  48 graphs",
              "",
              "Precision :  AMP FP16",
              "  (halves GPU memory, speeds training)",
              "",
              "Max epochs :  50",
              "Early stop patience :  10  (never triggered)"])

    add_card(slide, Inches(8.9), Inches(1.0), Inches(4.1), Inches(2.6),
             "Loss Function",
             ["Loss :  BCEWithLogitsLoss",
              "",
              "Class imbalance:  ~19 : 1",
              "  healthy vs. tumour superpixels",
              "",
              "Solution:  positive class weight",
              "  w+ = 9.0",
              "  Penalises missed tumour detections",
              "",
              "Dice smoothing ε = 1e-7"],
             title_color=ORANGE)

    add_card(slide, Inches(0.3), Inches(3.75), Inches(6.3), Inches(2.35),
             "5-Fold Cross-Validation",
             ["1,000 BraTS 2021 patients split at PATIENT level:",
              "",
              "  Each fold:   720 train  /  80 val  /  200 test",
              "",
              "5 models trained independently on 5 splits",
              "",
              "Ensemble:  soft voting — average sigmoid probabilities",
              "Threshold: τ = 0.5"],
             title_color=GREEN)

    add_card(slide, Inches(6.85), Inches(3.75), Inches(6.15), Inches(2.35),
             "Hardware (Consumer-Grade Only)",
             ["CPU :   Intel Core i7-10700",
              "RAM :   32 GB",
              "GPU :   NVIDIA RTX 2060  —  6 GB VRAM",
              "",
              "Python 3.12  ·  PyTorch 2.8.0  (CUDA 12.8)",
              "PyTorch Geometric 2.6.1",
              "scikit-image 0.25.2  ·  NumPy 2.3.3"],
             title_color=LIGHT_GRAY)

    add_text(slide,
             "Entire training pipeline runs on hardware available in most university labs.",
             Inches(0.3), Inches(6.25), Inches(12.7), Inches(0.45),
             font_size=Pt(12), color=LIGHT_GRAY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CV RESULTS + HELD-OUT (merged)
# ══════════════════════════════════════════════════════════════════════════════
def slide_results_cv(prs):
    slide = blank_slide(prs)
    add_header(slide, 10, "Experimental Results",
               "5-Fold Cross-Validation  +  Sealed Held-Out Set  (BraTS 2021)")

    # CV fold table
    add_rect(slide, Inches(0.3), Inches(1.0), Inches(5.8), Inches(4.4),
             CARD_BG, ACCENT_BLUE, line_width=Pt(1.2))
    add_text(slide, "5-Fold Cross-Validation  (1,000 patients)",
             Inches(0.45), Inches(1.06), Inches(5.5), Inches(0.38),
             font_size=Pt(14), bold=True, color=ACCENT_TEAL)

    folds = [
        ("Fold 0", "88.72%"),
        ("Fold 1", "90.48%"),
        ("Fold 2", "90.31%"),
        ("Fold 3", "90.13%"),
        ("Fold 4", "90.47%"),
        ("Mean ± σ", "90.02% ± 0.66%"),
    ]
    header_colors = [LIGHT_GRAY] * 5 + [ACCENT_TEAL]
    for i, ((fold, dice), col) in enumerate(zip(folds, header_colors)):
        ry = Inches(1.55) + i * Inches(0.49)
        is_mean = i == 5
        bg = ACCENT_BLUE if is_mean else CARD_BG
        add_rect(slide, Inches(0.45), ry, Inches(5.5), Inches(0.45), bg)
        add_text(slide, fold, Inches(0.55), ry + Inches(0.06), Inches(2.5), Inches(0.35),
                 font_size=Pt(13), bold=is_mean, color=col)
        add_text(slide, dice, Inches(3.4), ry + Inches(0.06), Inches(2.5), Inches(0.35),
                 font_size=Pt(13), bold=is_mean, color=col, align=PP_ALIGN.RIGHT)

    add_text(slide, "Stable across folds — σ = 0.66pp",
             Inches(0.45), Inches(4.72), Inches(5.5), Inches(0.35),
             font_size=Pt(12), color=LIGHT_GRAY, italic=True)

    # Held-out results
    add_rect(slide, Inches(6.4), Inches(1.0), Inches(6.6), Inches(4.4),
             CARD_BG, ACCENT_TEAL, line_width=Pt(1.5))
    add_text(slide, "Ensemble on Sealed Held-Out Set  (251 patients)",
             Inches(6.55), Inches(1.06), Inches(6.3), Inches(0.38),
             font_size=Pt(14), bold=True, color=ACCENT_TEAL)

    metrics = [
        ("Dice Score", "91.41%", ACCENT_TEAL, True),
        ("Accuracy", "99.14%", WHITE, False),
        ("Precision", "95.52%", WHITE, False),
        ("Sensitivity", "87.77%", WHITE, False),
        ("Specificity", "99.76%", WHITE, False),
    ]
    for i, (metric, value, col, is_primary) in enumerate(metrics):
        ry = Inches(1.55) + i * Inches(0.57)
        bg = ACCENT_BLUE if is_primary else CARD_BG
        add_rect(slide, Inches(6.55), ry, Inches(6.3), Inches(0.52), bg)
        add_text(slide, metric, Inches(6.65), ry + Inches(0.08), Inches(3.5), Inches(0.38),
                 font_size=Pt(13), bold=is_primary, color=col)
        add_text(slide, value, Inches(10.3), ry + Inches(0.08), Inches(2.4), Inches(0.38),
                 font_size=Pt(14 if is_primary else 13), bold=is_primary, color=col,
                 align=PP_ALIGN.RIGHT)

    # Highlight Dice big
    add_rect(slide, Inches(6.4), Inches(4.55), Inches(6.6), Inches(0.65),
             ACCENT_TEAL, None)
    add_text(slide, "Final Dice: 91.41%  on completely unseen 251 patients",
             Inches(6.55), Inches(4.6), Inches(6.3), Inches(0.55),
             font_size=Pt(15), bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

    add_text(slide,
             "Ensemble of 5 folds (soft voting, τ = 0.5) outperforms any single fold — confirms that cross-validation adds value.",
             Inches(0.3), Inches(5.55), Inches(12.7), Inches(0.45),
             font_size=Pt(12), color=LIGHT_GRAY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — EFFICIENCY ADVANTAGE
# ══════════════════════════════════════════════════════════════════════════════
def slide_efficiency(prs):
    slide = blank_slide(prs)
    add_header(slide, 11, "Efficiency Advantage",
               "GNN (ours) vs. U-Net — same hardware, same dataset")

    # Header row
    headers = ["Metric", "GNN (Ours)", "U-Net", "Improvement"]
    col_widths = [Inches(3.5), Inches(2.8), Inches(2.8), Inches(3.2)]
    col_starts = [Inches(0.3)]
    for cw in col_widths[:-1]:
        col_starts.append(col_starts[-1] + cw)

    header_y = Inches(1.0)
    add_rect(slide, Inches(0.3), header_y, sum(col_widths), Inches(0.48), ACCENT_BLUE)
    for ci, (hdr, cx, cw) in enumerate(zip(headers, col_starts, col_widths)):
        add_text(slide, hdr, cx + Inches(0.05), header_y + Inches(0.07),
                 cw - Inches(0.1), Inches(0.35),
                 font_size=Pt(13), bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

    rows = [
        ("Dice Score",              "91.41%",    "87.84%",    "+3.57 pp",     GREEN,   True),
        ("End-to-end Inference",    "1,732 ms",  "10,160 ms", "5.9× faster",  ACCENT_TEAL, False),
        ("GNN-only Inference",      "75.4 ms",   "—",         "—",            LIGHT_GRAY, False),
        ("Peak GPU Memory",         "11 MB",     "2,500 MB",  "227× less",    ACCENT_TEAL, False),
        ("Parameters",              "439,041",   "69,146,113","157× fewer",   ACCENT_TEAL, False),
        ("Model Storage",           "1.7 MB",    "264 MB",    "157× smaller", ACCENT_TEAL, False),
    ]

    for ri, (metric, ours, unet, impr, impr_col, highlight) in enumerate(rows):
        ry = Inches(1.52) + ri * Inches(0.72)
        bg = RGBColor(0x12, 0x26, 0x38) if ri % 2 == 0 else CARD_BG
        add_rect(slide, Inches(0.3), ry, sum(col_widths), Inches(0.68), bg)
        vals = [metric, ours, unet, impr]
        colors = [WHITE, GREEN if highlight else WHITE, RED if highlight else LIGHT_GRAY, impr_col]
        for ci, (val, cx, cw, vc) in enumerate(zip(vals, col_starts, col_widths, colors)):
            add_text(slide, val, cx + Inches(0.05), ry + Inches(0.15),
                     cw - Inches(0.1), Inches(0.42),
                     font_size=Pt(13), bold=(ci == 3),
                     color=vc,
                     align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

    add_rect(slide, Inches(0.3), Inches(5.9), Inches(12.7), Inches(0.55),
             CARD_BG, GREEN, line_width=Pt(1.5))
    add_text(slide,
             "Our GNN is MORE accurate than U-Net (+3.57pp Dice) while using 227× less memory, 5.9× faster, with 157× fewer parameters.",
             Inches(0.5), Inches(5.97), Inches(12.3), Inches(0.45),
             font_size=Pt(13), bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    add_text(slide,
             "75.4 ms GNN-only inference = real-time capable. The 1,732 ms total includes superpixel preprocessing.",
             Inches(0.3), Inches(6.58), Inches(12.7), Inches(0.42),
             font_size=Pt(12), color=LIGHT_GRAY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — BRATS 2023 ZERO-SHOT
# ══════════════════════════════════════════════════════════════════════════════
def slide_zeroshot(prs):
    slide = blank_slide(prs)
    add_header(slide, 12, "Zero-Shot Generalisation",
               "BraTS 2023 — no retraining, different acquisition protocols")

    # Explanation
    add_card(slide, Inches(0.3), Inches(1.0), Inches(5.5), Inches(1.7),
             "What is Zero-Shot Transfer?",
             ["Model trained ONLY on BraTS 2021",
              "Tested directly on BraTS 2023 — no fine-tuning",
              "",
              "BraTS 2023 uses different scanners & protocols",
              "→ This tests real-world robustness"],
             title_color=ACCENT_TEAL, body_size=Pt(12))

    add_card(slide, Inches(6.1), Inches(1.0), Inches(6.9), Inches(1.7),
             "Why GraphSAGE Generalises",
             ["Inductive framework — learns an aggregation function",
              "not a fixed adjacency matrix",
              "",
              "Superpixel features are scanner-agnostic relative statistics",
              "→ mean/std intensities are invariant to protocol shifts"],
             title_color=GREEN, body_size=Pt(12))

    # Comparison table
    headers = ["Metric", "BraTS 2021", "BraTS 2023", "Change"]
    col_widths = [Inches(3.2), Inches(2.9), Inches(2.9), Inches(3.0)]
    col_starts = [Inches(0.3)]
    for cw in col_widths[:-1]:
        col_starts.append(col_starts[-1] + cw)

    header_y = Inches(2.9)
    add_rect(slide, Inches(0.3), header_y, sum(col_widths), Inches(0.48), ACCENT_BLUE)
    for ci, (hdr, cx, cw) in enumerate(zip(headers, col_starts, col_widths)):
        add_text(slide, hdr, cx + Inches(0.05), header_y + Inches(0.07),
                 cw - Inches(0.1), Inches(0.35),
                 font_size=Pt(13), bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

    rows = [
        ("Dice Score",   "91.41%", "89.40%", "−1.01 pp", LIGHT_GRAY),
        ("Accuracy",     "99.14%", "98.85%", "−0.29 pp", LIGHT_GRAY),
        ("Sensitivity",  "87.77%", "90.69%", "+2.92 pp ↑", GREEN),
        ("Specificity",  "99.76%", "99.45%", "−0.31 pp", LIGHT_GRAY),
        ("Precision",    "95.52%", "92.46%", "−3.06 pp", LIGHT_GRAY),
    ]
    for ri, (metric, v21, v23, change, change_col) in enumerate(rows):
        ry = Inches(3.42) + ri * Inches(0.52)
        bg = RGBColor(0x12, 0x26, 0x38) if ri % 2 == 0 else CARD_BG
        add_rect(slide, Inches(0.3), ry, sum(col_widths), Inches(0.48), bg)
        for ci, (val, cx, cw, vc) in enumerate(zip(
                [metric, v21, v23, change],
                col_starts, col_widths,
                [WHITE, WHITE, WHITE, change_col])):
            add_text(slide, val, cx + Inches(0.05), ry + Inches(0.08),
                     cw - Inches(0.1), Inches(0.35),
                     font_size=Pt(13), bold=(ci == 3 and change_col == GREEN),
                     color=vc,
                     align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

    add_rect(slide, Inches(0.3), Inches(6.1), Inches(12.7), Inches(0.55),
             CARD_BG, GREEN, line_width=Pt(1.5))
    add_text(slide,
             "Sensitivity IMPROVED +2.92pp on unseen data — model detects tumour better on BraTS 2023 than BraTS 2021.",
             Inches(0.5), Inches(6.17), Inches(12.3), Inches(0.42),
             font_size=Pt(13), bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    add_text(slide,
             "Modest −1.01pp Dice drop with zero retraining confirms strong cross-dataset generalisation.",
             Inches(0.3), Inches(6.75), Inches(12.7), Inches(0.38),
             font_size=Pt(12), color=LIGHT_GRAY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — LIMITATIONS + FUTURE WORK (merged)
# ══════════════════════════════════════════════════════════════════════════════
def slide_limitations_future(prs):
    slide = blank_slide(prs)
    add_header(slide, 13, "Limitations  &  Future Work",
               "Honest assessment + next research directions")

    # Limitations
    limitations = [
        ("T1CE Dependency",
         "SLIC on T1CE fails for non-enhancing tumours — 3 patients had Dice ≈ 0 "
         "(BraTS2021_01405, 01366, 01407). No T1CE signal = no usable superpixels.",
         RED),
        ("Binary Segmentation Only",
         "Predicts Whole Tumour only. Does not separate sub-regions: "
         "GD-enhancing tumour, necrotic core, or peritumoral oedema.",
         ORANGE),
        ("Fixed Graph Topology",
         "Graph structure is fixed at inference time — no dynamic refinement based on "
         "prediction confidence or iterative superpixel adjustment.",
         YELLOW),
        ("Adult Glioma Only",
         "Trained exclusively on adult glioma cases (BraTS). Performance on paediatric "
         "tumours or rare tumour types is unknown.",
         LIGHT_GRAY),
    ]

    for i, (title, body, color) in enumerate(limitations):
        col = i % 2
        row = i // 2
        x = Inches(0.3) + col * Inches(6.4)
        y = Inches(1.05) + row * Inches(2.0)
        add_card(slide, x, y, Inches(6.15), Inches(1.8), title, [body],
                 title_color=color, title_size=Pt(13), body_size=Pt(11))

    # Future work
    add_rect(slide, Inches(0.3), Inches(5.2), Inches(12.7), Inches(0.38),
             ACCENT_BLUE)
    add_text(slide, "Future Work",
             Inches(0.5), Inches(5.24), Inches(3.0), Inches(0.32),
             font_size=Pt(13), bold=True, color=WHITE)

    future = [
        "Multi-class segmentation (WT, TC, ET sub-regions)",
        "Dynamic graph refinement based on prediction confidence",
        "Extend to paediatric BraTS datasets",
        "Integrate radiomics / clinical metadata as node features",
        "Real-time deployment on edge hardware (Jetson Nano / Raspberry Pi)",
    ]
    add_text(slide, "   ·   ".join(future),
             Inches(0.3), Inches(5.68), Inches(12.7), Inches(0.8),
             font_size=Pt(11), color=LIGHT_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — CONCLUSIONS
# ══════════════════════════════════════════════════════════════════════════════
def slide_conclusion(prs):
    slide = blank_slide(prs)
    add_header(slide, 14, "Conclusions",
               "What we demonstrated — four measurable achievements")

    conclusions = [
        ("01",
         "91.41% Dice — Surpasses U-Net",
         "GraphSAGE on SLIC superpixels achieves 91.41% Dice on the sealed BraTS 2021 held-out "
         "set — beating U-Net baseline (87.84%) by +3.57 percentage points.",
         GREEN),
        ("02",
         "Radical Efficiency Gains",
         "5.9× faster inference  ·  227× less GPU memory  ·  157× fewer parameters.\n"
         "Full pipeline runs on a 6 GB consumer GPU — no HPC cluster required.",
         ACCENT_TEAL),
        ("03",
         "Zero-Shot Generalisation Confirmed",
         "Without any retraining, 89.40% Dice on BraTS 2023 — a different acquisition protocol.\n"
         "Sensitivity actually improved +2.92pp, validating GraphSAGE's inductive design.",
         YELLOW),
        ("04",
         "Near-SOTA at Fraction of Compute",
         "SOTA CNNs achieve 92–94% Dice but need 16–32 GB VRAM and 100M+ params.\n"
         "Our GNN is within 1–3pp of SOTA at 439K params and 11 MB peak memory.",
         ORANGE),
    ]

    for i, (num, title, body, color) in enumerate(conclusions):
        col = i % 2
        row = i // 2
        x = Inches(0.3) + col * Inches(6.5)
        y = Inches(1.05) + row * Inches(2.7)
        w = Inches(6.2)
        h = Inches(2.45)
        add_rect(slide, x, y, w, h, CARD_BG, color, line_width=Pt(2))
        add_text(slide, num,
                 x + Inches(0.12), y + Inches(0.08), Inches(0.65), Inches(0.65),
                 font_size=Pt(28), bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, title,
                 x + Inches(0.85), y + Inches(0.1), w - Inches(1.0), Inches(0.45),
                 font_size=Pt(15), bold=True, color=WHITE)
        add_text(slide, body,
                 x + Inches(0.12), y + Inches(0.65), w - Inches(0.25), h - Inches(0.75),
                 font_size=Pt(12), color=LIGHT_GRAY)

    add_rect(slide, Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.55),
             ACCENT_TEAL, None)
    add_text(slide,
             "GraphSAGE + SLIC superpixels = near-SOTA accuracy with clinical-grade efficiency on consumer hardware.",
             Inches(0.5), Inches(6.67), Inches(12.3), Inches(0.42),
             font_size=Pt(14), bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — THANK YOU / Q&A
# ══════════════════════════════════════════════════════════════════════════════
def slide_thankyou(prs):
    slide = blank_slide(prs)
    fill_bg(slide, BG_DARK)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), ACCENT_TEAL)
    add_rect(slide, 0, Inches(7.42), SLIDE_W, Inches(0.08), ACCENT_TEAL)

    # Central block
    add_rect(slide, Inches(1.5), Inches(1.5), Inches(10.33), Inches(4.5),
             CARD_BG, ACCENT_TEAL, line_width=Pt(2))
    add_text(slide, "Thank You",
             Inches(1.6), Inches(1.7), Inches(10.13), Inches(1.0),
             font_size=Pt(52), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "Questions & Discussion",
             Inches(1.6), Inches(2.7), Inches(10.13), Inches(0.65),
             font_size=Pt(26), color=ACCENT_TEAL, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(1.8), Inches(3.55), Inches(9.73), Inches(0.04), ACCENT_BLUE)

    # Key numbers reminder
    stats = [
        ("91.41%", "Dice Score"),
        ("5.9×", "Faster"),
        ("227×", "Less Memory"),
        ("89.40%", "Zero-Shot"),
    ]
    sw = Inches(2.3)
    sx = Inches(2.1)
    for i, (num, label) in enumerate(stats):
        x = sx + i * sw
        add_text(slide, num, x, Inches(3.72), sw, Inches(0.55),
                 font_size=Pt(22), bold=True, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)
        add_text(slide, label, x, Inches(4.28), sw, Inches(0.3),
                 font_size=Pt(12), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    add_text(slide,
             "Sakib Khan  ·  Rifa Sanjida  ·  Kishor Kumar Das  ·  Md. Mahamudul Hasan  ·  Md. Minhajur Rahman",
             Inches(0.5), Inches(5.35), Inches(12.33), Inches(0.4),
             font_size=Pt(13), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, "Supervisor: Mr. Shamim Ahmed  ·  BUBT  ·  2026",
             Inches(0.5), Inches(5.78), Inches(12.33), Inches(0.38),
             font_size=Pt(12), color=YELLOW, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════
if __name__ == "__main__":
    prs = new_prs()

    slide_title(prs)               # 1 — Title + Team
    slide_problem(prs)             # 2 — Problem Statement
    slide_objectives(prs)          # 3 — Research Objectives
    slide_dataset(prs)             # 4 — Dataset
    slide_methodology_overview(prs)# 5 — Methodology Overview
    slide_slic(prs)                # 6 — SLIC Superpixels
    slide_graph_construction(prs)  # 7 — Graph Construction
    slide_features_graphsage(prs)  # 8 — Node Features + GraphSAGE
    slide_training(prs)            # 9 — Training Protocol
    slide_results_cv(prs)          # 10 — CV + Held-Out Results
    slide_efficiency(prs)          # 11 — Efficiency Advantage
    slide_zeroshot(prs)            # 12 — BraTS 2023 Zero-Shot
    slide_limitations_future(prs)  # 13 — Limitations + Future Work
    slide_conclusion(prs)          # 14 — Conclusions
    slide_thankyou(prs)            # 15 — Thank You / Q&A

    out = "/mnt/bigdata/capstone/brats_gnn_segmentation/paperWriting/compressed_15_slides.pptx"
    prs.save(out)
    print(f"Saved → {out}")
    print(f"Total slides: {len(prs.slides)}")
