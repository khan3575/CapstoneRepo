"""
Build Capstone_Presentation_v2.pptx — 27 slides aligned to PRESENTATION_SPEECH.md

KEY APPROACH: copy source PPTX first (preserves embedded fonts, theme, media),
then reorder / delete / add slides entirely within that copy.
This avoids the font-stripping bug that broke Material Symbols icons in v1.

Final 27-slide order (source indices, 0-based):
  pos  1  → src[0]   Title
  pos  2  → src[3]   Introduction
  pos  3  → src[4]   Problem Statement
  pos  4  → NEW      Research Gap
  pos  5  → NEW      Proposed Approach
  pos  6  → src[5]   Research Objectives
  pos  7  → src[9]   Dataset: BraTS
  pos  8  → NEW      MRI Modalities
  pos  9  → src[10]  Methodology Overview (diagram)
  pos 10  → src[11]  SLIC Superpixels
  pos 11  → src[12]  Graph Construction
  pos 12  → src[13]  Node Features
  pos 13  → NEW      GraphSAGE Architecture (visualization image)
  pos 14  → src[15]  Training Protocol
  pos 15  → src[17]  Evaluation Metrics
  pos 16  → src[18]  Cross-Validation Results
  pos 17  → src[19]  Held-Out Test Results  (BraTS-2023 footnote removed)
  pos 18  → src[20]  Comparison with other papers
  pos 19  → src[21]  Efficiency Comparison
  pos 20  → src[22]  BraTS 2023 Generalisation
  pos 21  → src[23]  Failure Cases
  pos 22  → src[24]  Ablation Study
  pos 23  → src[26]  Limitations
  pos 24  → src[27]  Future Work
  pos 25  → src[28]  Conclusions
  pos 26  → src[29]  References
  pos 27  → src[30]  Thank You / Q&A

Dropped (7 slides): src[1] Presented By, src[2] TOC, src[6] Key Contributions,
  src[7] Background & Motivation, src[8] Related Works,
  src[16] Experimental Setup, src[25] Key Project Achievements
"""

import shutil
import os
from io import BytesIO
from lxml import etree
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# ─────────────────────────────────────────────
# Paths
# ─────────────────────────────────────────────
BASE    = os.path.dirname(__file__)
SRC     = os.path.join(BASE, "final Presentation.pptx")
OUT     = os.path.join(BASE, "Capstone_Presentation_v2.pptx")

# ─────────────────────────────────────────────
# Theme colours (matching existing dark slides)
# ─────────────────────────────────────────────
C_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
C_CYAN        = RGBColor(0x00, 0xB4, 0xD8)
C_LIGHT_GRAY  = RGBColor(0xCC, 0xD6, 0xE0)
C_MID_FILL    = RGBColor(0x1A, 0x2E, 0x44)
C_DEEP_FILL   = RGBColor(0x0A, 0x1E, 0x30)
C_CYAN_DIM    = RGBColor(0x00, 0x70, 0x90)

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)


# ─────────────────────────────────────────────
# Low-level helpers
# ─────────────────────────────────────────────

def _add_background(slide, img_blob):
    """Insert a full-bleed image at the back of the slide."""
    pic = slide.shapes.add_picture(BytesIO(img_blob), 0, 0, SLIDE_W, SLIDE_H)
    sp  = pic._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)


def _txb(slide, left, top, w, h, text,
         font="Calibri", size=12, bold=False, color=C_LIGHT_GRAY,
         bg=None, align=PP_ALIGN.LEFT, wrap=True):
    """Add a text box with a single run."""
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    if bg:
        tb.fill.solid()
        tb.fill.fore_color.rgb = bg
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = font
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return tb


def _rule(slide, y=Inches(1.20), width=Inches(9.16)):
    r = slide.shapes.add_shape(1, Inches(0.42), y, width, Emu(28800))
    r.fill.solid()
    r.fill.fore_color.rgb = C_CYAN_DIM
    r.line.fill.background()


def _header(slide, title, subtitle=None):
    _txb(slide, Inches(0.42), Inches(0.26), Inches(9.16), Inches(0.55),
         title, size=28, bold=True, color=C_WHITE)
    if subtitle:
        _txb(slide, Inches(0.42), Inches(0.82), Inches(9.16), Inches(0.30),
             subtitle, size=12, color=C_CYAN)


def _card(slide, left, top, w, h, label, body,
          lbl_size=11, body_size=9.5):
    """Dark panel with cyan label + light-gray body text."""
    pan = slide.shapes.add_shape(1, left, top, w, h)
    pan.fill.solid()
    pan.fill.fore_color.rgb = C_MID_FILL
    pan.line.color.rgb      = C_CYAN_DIM
    pan.line.width          = Pt(0.75)
    _txb(slide, left + Inches(0.13), top + Inches(0.09),
         w - Inches(0.26), Inches(0.28),
         label, size=lbl_size, bold=True, color=C_CYAN)
    _txb(slide, left + Inches(0.13), top + Inches(0.39),
         w - Inches(0.26), h - Inches(0.48),
         body, size=body_size, color=C_LIGHT_GRAY, wrap=True)


# ─────────────────────────────────────────────
# New-slide builders
# ─────────────────────────────────────────────

def _new_slide(prs):
    """Add a blank slide to the presentation and return it."""
    # Use the layout that the dark-theme content slides use (Statement layout)
    # Any layout works since we cover it with a full-bleed background image.
    layout = prs.slide_layouts[0]
    return prs.slides.add_slide(layout)


def build_research_gap(prs, bg_blob):
    sl = _new_slide(prs)
    _add_background(sl, bg_blob)
    _header(sl, "Research Gap",
            "What the existing literature was missing")
    _rule(sl, y=Inches(1.22))

    W, H = Inches(4.60), Inches(1.82)
    cards = [
        (Inches(0.42), Inches(1.30),
         "Accuracy-Only Focus",
         "Nearly all SOTA methods (nnU-Net, SwinUNETR, TransBTS) optimise solely for "
         "Dice score. Efficiency metrics — memory, speed, model size — receive almost "
         "no systematic attention."),
        (Inches(5.20), Inches(1.30),
         "Hardware Inaccessibility",
         "Leading CNN/Transformer models require 16–32 GB of GPU VRAM. Most hospitals "
         "worldwide cannot afford this infrastructure, blocking real-world adoption."),
        (Inches(0.42), Inches(3.22),
         "Massive Parameter Counts",
         "nnU-Net: 31 M  ·  Swin-UNETR: 62 M  ·  TransBTS: 33 M. "
         "Hundreds of millions of parameters with no controlled efficiency benchmarking "
         "against the hardware constraints of clinical settings."),
        (Inches(5.20), Inches(3.22),
         "GNN Work Lacked Scale",
         "Prior GNN-based segmentation studies (Saueressig et al., Patel et al.) used "
         "small cohorts and provided no direct hardware-matched comparison against "
         "standard CNN baselines."),
    ]
    for lft, tp, lbl, bdy in cards:
        _card(sl, lft, tp, W, H, lbl, bdy)
    return sl


def build_proposed_approach(prs, bg_blob):
    sl = _new_slide(prs)
    _add_background(sl, bg_blob)
    _header(sl, "Proposed Approach",
            "Graph Neural Network Pipeline for Brain Tumour Segmentation")
    _rule(sl, y=Inches(1.22))

    # Core-idea banner
    ban = sl.shapes.add_shape(1, Inches(0.42), Inches(1.30),
                               Inches(9.16), Inches(0.70))
    ban.fill.solid()
    ban.fill.fore_color.rgb = RGBColor(0x07, 0x1E, 0x32)
    ban.line.color.rgb      = C_CYAN
    ban.line.width          = Pt(1)
    _txb(sl, Inches(0.58), Inches(1.33), Inches(8.82), Inches(0.64),
         "Key Idea: Convert 3D MRI volumes into compact superpixel graphs, then "
         "apply a lightweight GNN for node-level binary tumour classification.",
         size=11, color=C_WHITE, wrap=True)

    # Four step cards
    W, H = Inches(4.60), Inches(1.48)
    steps = [
        (Inches(0.42), Inches(2.10),
         "Step 1 — SLIC Superpixels",
         "Each 2D axial MRI slice is segmented into ~46 compact superpixels using "
         "SLIC on the T1CE channel, achieving 808× per-slice spatial compression."),
        (Inches(5.20), Inches(2.10),
         "Step 2 — Graph Construction",
         "Two consecutive slices are merged into one graph: ~92 nodes, ~180 edges "
         "(intra-slice RAG + inter-slice k-NN with k=3, overlap >10% or ≤10 mm)."),
        (Inches(0.42), Inches(3.68),
         "Step 3 — GraphSAGE Classifier",
         "5-layer inductive GNN, 256 hidden dims, 439 K parameters. "
         "Mean-pooling aggregation. Binary node classification: tumour vs. healthy."),
        (Inches(5.20), Inches(3.68),
         "Step 4 — Soft-Voting Ensemble",
         "5-fold cross-validation models combine predictions via sigmoid averaging "
         "(threshold = 0.5) for the final patient-level segmentation mask."),
    ]
    for lft, tp, lbl, bdy in steps:
        _card(sl, lft, tp, W, H, lbl, bdy)

    # Stat strip
    strip = sl.shapes.add_shape(1, Inches(0.42), Inches(5.24),
                                 Inches(9.16), Inches(0.27))
    strip.fill.solid()
    strip.fill.fore_color.rgb = RGBColor(0x00, 0x2A, 0x40)
    strip.line.fill.background()
    _txb(sl, Inches(0.55), Inches(5.25), Inches(8.90), Inches(0.26),
         "2,284× spatial compression  ·  5.9× faster inference  ·  "
         "227× less GPU memory  ·  157× fewer parameters",
         size=9, color=C_CYAN)
    return sl


def build_graphsage_architecture(prs, bg_blob):
    """GraphSAGE Architecture slide — full-width visualization image + header."""
    sl = _new_slide(prs)
    _add_background(sl, bg_blob)
    _header(sl, "GraphSAGE Architecture",
            "Graph Sample and Aggregate — 5-layer inductive GNN for node-level tumour classification")
    _rule(sl, y=Inches(1.22))

    img_path = os.path.join(BASE, "graphsage_architecture.png")
    if not os.path.exists(img_path):
        raise FileNotFoundError(f"Missing image: {img_path}\n"
                                "Run gen_graphsage_visual.py first.")

    # Full-width image, centred vertically in the remaining space
    img_w = Inches(9.40)
    img_h = Inches(3.95)
    img_l = (SLIDE_W - img_w) / 2
    img_t = Inches(1.36)

    with open(img_path, "rb") as f:
        sl.shapes.add_picture(f, img_l, img_t, img_w, img_h)

    # Bottom key-stats strip
    strip = sl.shapes.add_shape(1, Inches(0.42), Inches(5.28),
                                 Inches(9.16), Inches(0.24))
    strip.fill.solid()
    strip.fill.fore_color.rgb = RGBColor(0x00, 0x1E, 0x30)
    strip.line.fill.background()
    _txb(sl, Inches(0.55), Inches(5.28), Inches(8.90), Inches(0.24),
         "5 layers  ·  256 hidden dims  ·  64 output dims  ·  "
         "Mean-pool aggregation  ·  439,041 parameters  ·  ReLU + BatchNorm",
         size=8.5, color=C_CYAN)
    return sl


def build_mri_modalities(prs, bg_blob):
    sl = _new_slide(prs)
    _add_background(sl, bg_blob)
    _header(sl, "MRI Modalities",
            "Four complementary views of the brain — each modality reveals different tissue properties")
    _rule(sl, y=Inches(1.22))

    W, H = Inches(4.50), Inches(1.88)
    modalities = [
        (Inches(0.42), Inches(1.32),
         "T1", "Basic anatomical structure",
         "Provides the standard anatomical reference of the brain — distinguishes grey "
         "matter, white matter, and CSF. Used as a baseline anatomical map by clinicians."),
        (Inches(5.08), Inches(1.32),
         "T1CE  (T1 Contrast Enhanced)", "Active tumour — blood-brain barrier breakdown",
         "Gadolinium contrast agent highlights regions where the blood-brain barrier is "
         "disrupted. Active tumour appears bright. This channel provides the clearest tumour "
         "boundary signal and is the basis for SLIC superpixel construction."),
        (Inches(0.42), Inches(3.30),
         "T2", "Water content — oedema and swelling",
         "Water-sensitive sequence that makes oedema (swelling) around the tumour very "
         "visible. Tumour infiltration zone and peritumoral oedema appear hyperintense (bright)."),
        (Inches(5.08), Inches(3.30),
         "FLAIR  (Fluid Attenuated Inversion Recovery)", "Abnormal tissue — CSF suppressed",
         "Suppresses normal cerebrospinal fluid signal, making abnormal infiltrated tissue "
         "at the tumour margin much clearer. Particularly sensitive to oedema."),
    ]
    for lft, tp, modal, tagline, detail in modalities:
        pan = sl.shapes.add_shape(1, lft, tp, W, H)
        pan.fill.solid()
        pan.fill.fore_color.rgb = C_MID_FILL
        pan.line.color.rgb      = C_CYAN_DIM
        pan.line.width          = Pt(0.75)
        _txb(sl, lft + Inches(0.13), tp + Inches(0.08),
             W - Inches(0.26), Inches(0.28),
             modal, size=12, bold=True, color=C_CYAN)
        _txb(sl, lft + Inches(0.13), tp + Inches(0.39),
             W - Inches(0.26), Inches(0.23),
             tagline, size=9.5, bold=True, color=C_WHITE)
        _txb(sl, lft + Inches(0.13), tp + Inches(0.64),
             W - Inches(0.26), H - Inches(0.72),
             detail, size=9, color=C_LIGHT_GRAY, wrap=True)

    _txb(sl, Inches(0.42), Inches(5.30), Inches(9.16), Inches(0.22),
         "All four modalities are fused as node features — giving the model a "
         "multi-contrast tissue signature for each superpixel.",
         size=8.5, color=C_CYAN)
    return sl


# ─────────────────────────────────────────────
# Slide-list reordering helper
# ─────────────────────────────────────────────

def reorder_slides(prs, desired_indices):
    """
    Reorder the presentation slides so that only the slides in
    desired_indices appear, in that order.
    desired_indices: 0-based indices into the CURRENT slide list.
    Slides not mentioned are dropped (their rels are removed).
    """
    sldIdLst = prs.slides._sldIdLst
    all_sldId = list(sldIdLst)        # snapshot before any mutation

    # Identify rIds to keep vs. drop
    keep_rIds = {all_sldId[i].get(qn("r:id")) for i in desired_indices}
    drop_rIds = [elem.get(qn("r:id")) for elem in all_sldId
                 if elem.get(qn("r:id")) not in keep_rIds]

    # Clear the entire list
    for elem in all_sldId:
        sldIdLst.remove(elem)

    # Re-add in desired order
    for i in desired_indices:
        sldIdLst.append(all_sldId[i])

    # Remove relationships for dropped slides (access underlying XML directly)
    rels_element = prs.part._element.getparent()   # may not exist; use rels XML
    # Access the .rels part via the package
    try:
        rel_part = prs.part.rels
        for rId in drop_rIds:
            # _Relationships stores items in _rels dict internally
            if hasattr(rel_part, '_rels') and rId in rel_part._rels:
                del rel_part._rels[rId]
            elif hasattr(rel_part, '_relsByRId') and rId in rel_part._relsByRId:
                del rel_part._relsByRId[rId]
    except Exception as e:
        print(f"  Note: could not remove orphaned slide rels ({e}) — harmless.")


# ─────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────

def main():
    # ── 1. Copy source → output (preserves fonts, theme, media) ──
    shutil.copy(SRC, OUT)
    prs = Presentation(OUT)
    print(f"Opened source: {len(prs.slides)} slides")

    # ── 2. Extract dark background blob ──
    bg_blob = None
    for shape in prs.slides[3].shapes:   # slide 4 (Introduction) = dark bg
        if shape.shape_type == 13:
            bg_blob = shape.image.blob
            break
    assert bg_blob, "Could not find background image"

    # ── 3. Fix slide 19 (Held-Out): remove stray BraTS-2023 footnote ──
    held_out = prs.slides[19]
    for shape in list(held_out.shapes):
        if shape.name == "Google Shape;1090;p105":
            shape._element.getparent().remove(shape._element)
            print("  Removed BraTS-2023 footnote from slide 19")
            break

    # ── 4. Add the 3 new slides (appended at indices 31, 32, 33) ──
    print("  Building Research Gap slide …")
    build_research_gap(prs, bg_blob)      # → index 31

    print("  Building Proposed Approach slide …")
    build_proposed_approach(prs, bg_blob) # → index 32

    print("  Building MRI Modalities slide …")
    build_mri_modalities(prs, bg_blob)    # → index 33

    print("  Building GraphSAGE Architecture slide …")
    build_graphsage_architecture(prs, bg_blob)  # → index 34

    print(f"  Slides after additions: {len(prs.slides)}")

    # ── 5. Strip Material Symbols icon-only subshapes ──────────────────────
    # Several source slides contain decorative GROUP subshapes whose ENTIRE
    # text is a Material Symbols ligature code ('timer', 'biotech', etc.).
    # Without the font installed, these render as huge default characters that
    # overflow the card and push all real content off-screen.
    # Fix: remove any sp/pic element inside a grpSp whose full text == an icon code.
    ICON_CODES = {
        # Problem Statement (src[4])
        'timer', 'storage', 'grid_view', 'science',
        # Research Objectives (src[5])
        'account_tree', 'psychology', 'speed', 'public',
        # Limitations (src[26])
        'layers', 'biotech', 'category', 'memory',
        # Future Works (src[27])
        'auto_graph', 'medical_services',
        # Conclusions (src[28])
        'workspace_premium', 'security', 'health_and_safety',
        # misc that may appear elsewhere
        'computer', 'compare', 'trending_up', 'verified',
        'groups', 'shield', 'local_hospital',
    }

    def _strip_icons(elem):
        """Recursively remove icon-only shapes from a spTree or grpSp element."""
        for child in list(elem):
            local = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if local == 'grpSp':
                _strip_icons(child)          # recurse into group
            elif local in ('sp', 'pic'):
                texts = child.xpath('.//*[local-name()="t"]/text()')
                full  = ''.join(texts).strip()
                if full in ICON_CODES:
                    elem.remove(child)

    # Apply to the slides that are affected (original 0-based indices)
    ICON_SLIDES = [4, 5, 26, 27, 28]   # Problem Stmt, Res.Obj, Limits, FutWork, Concl
    for idx in ICON_SLIDES:
        _strip_icons(prs.slides[idx].shapes._spTree)
    print(f"  Stripped icon shapes from slides: {ICON_SLIDES}")

    # ── 6. Reorder to the final 27-slide sequence ──
    # Indices reference the list AFTER step 4 (35 slides total, 0-34)
    DESIRED = [
        0,   # Title
        3,   # Introduction
        4,   # Problem Statement
        31,  # Research Gap        [NEW]
        32,  # Proposed Approach   [NEW]
        5,   # Research Objectives
        9,   # Dataset: BraTS
        33,  # MRI Modalities      [NEW]
        10,  # Methodology Overview
        11,  # SLIC Superpixels
        12,  # Graph Construction
        13,  # Node Features
        34,  # GraphSAGE Architecture [NEW — replaces src[14]]
        15,  # Training Protocol
        17,  # Evaluation Metrics
        18,  # Cross-Validation Results
        19,  # Held-Out Test Results (footnote removed)
        20,  # Comparison with other papers
        21,  # Efficiency Comparison
        22,  # BraTS 2023 Generalisation
        23,  # Failure Cases
        24,  # Ablation Study
        26,  # Limitations
        27,  # Future Work
        28,  # Conclusions
        29,  # References
        30,  # Thank You / Q&A
    ]
    assert len(DESIRED) == 27, f"Expected 27, got {len(DESIRED)}"

    reorder_slides(prs, DESIRED)
    print(f"  Slides after reorder: {len(prs.slides)}")

    # ── 7. Save ──
    prs.save(OUT)
    print(f"\nSaved → {OUT}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
